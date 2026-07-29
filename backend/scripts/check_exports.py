"""Build and parse every deterministic export format without external calls."""

from io import BytesIO

from openpyxl import load_workbook
from pptx import Presentation

from app.services.pitch_deck import PitchDeckGenerator
from app.utils.exports import export_financials_as_excel, export_plan_as_pdf


PLAN = {
    "title": "Prooflane Export Check",
    "idea": "An evidence-first purchasing assistant for independent restaurants.",
    "market_analysis": (
        "## Market Overview\n"
        "- Founder assumption TAM: $1 billion.\n"
        "## Target Customers\n"
        "- Independent restaurant operators."
    ),
    "business_model": (
        "## Value Proposition\n"
        "- Reduce purchasing waste.\n"
        "## Revenue Streams\n"
        "- Monthly subscription."
    ),
    "validation_strategy": (
        "## Viability Score: 64/100\n"
        "### Verdict: PROMISING\n"
        "### Validation Recommendations\n"
        "1. Interview 20 operators."
    ),
    "risks": "## Critical Risks\n- Supplier data quality is unproven.",
    "financials": "## Revenue Model\n- Subscription assumptions require validation.",
    "full_plan": (
        "# Business Plan: Prooflane Export Check\n"
        "## 1. Executive Summary\n"
        "An evidence-first purchasing assistant.\n"
        "## 3. Market Analysis\n"
        "A founder-estimated opportunity.\n"
        "## 4. Products & Services\n"
        "A purchasing recommendation workspace.\n"
        "## 5. Business Model\n"
        "Subscription revenue.\n"
        "## 6. Marketing & Sales Strategy\n"
        "Operator interviews and pilots.\n"
        "## 10. Risk Analysis\n"
        "Data quality must be tested.\n"
        "## 11. Funding Requirements\n"
        "Validation capital.\n"
        "## 12. Implementation Timeline\n"
        "Run interviews, then a paid pilot."
    ),
    "viability_score": 64,
    "intelligence": {
        "identity": {
            "name": "Prooflane Export Check",
            "subtitle": "Restaurant purchasing intelligence",
            "one_liner": "Reduce waste with evidence-backed ordering.",
        },
        "adjusted_score": {
            "raw_score": 64,
            "adjusted_score": 52,
            "evidence_confidence": 30,
        },
        "financial_metrics": {},
        "contradiction_issues": [],
        "consistency_issues": [],
    },
}

METRICS = {
    "year_1_revenue": {
        "value": 100_000,
        "display": "$100K",
        "unit": "USD",
        "basis": "projection",
        "confidence": 30,
        "assumption": True,
    },
    "use_of_funds": [
        {"category": "Validation", "percentage": 100, "amount": 250_000}
    ],
    "assumptions": [
        {
            "name": "CAC",
            "value": "$30",
            "rationale": "Founder estimate",
            "confidence": 20,
            "impact": "high",
            "validation_method": "Paid acquisition test",
        }
    ],
}


if __name__ == "__main__":
    pdf = export_plan_as_pdf(PLAN)
    assert pdf.startswith(b"%PDF")
    assert len(pdf) > 5_000

    pptx = PitchDeckGenerator().generate(PLAN)
    presentation = Presentation(BytesIO(pptx))
    assert len(presentation.slides) >= 10

    xlsx = export_financials_as_excel(PLAN["financials"], METRICS)
    workbook = load_workbook(BytesIO(xlsx))
    assert {"Key Metrics", "Use of Funds", "Assumptions"}.issubset(workbook.sheetnames)

    print(
        {
            "status": "passed",
            "pdf_bytes": len(pdf),
            "pitch_slides": len(presentation.slides),
            "workbook_sheets": workbook.sheetnames,
        }
    )
