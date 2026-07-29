"""Deterministic, investor-ready PowerPoint pitch deck generation."""
from __future__ import annotations

import re
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


INK = RGBColor(11, 18, 32)
INK_2 = RGBColor(22, 32, 51)
WHITE = RGBColor(250, 250, 248)
MUTED = RGBColor(166, 178, 194)
TEAL = RGBColor(45, 212, 191)
CORAL = RGBColor(249, 115, 96)
BLUE = RGBColor(82, 130, 255)
CARD = RGBColor(25, 37, 58)
LINE = RGBColor(53, 68, 91)


def _clean(value: str) -> str:
    value = re.sub(r"\*\*(.+?)\*\*", r"\1", value or "")
    value = re.sub(r"(?<!\*)\*([^*]+?)\*(?!\*)", r"\1", value)
    value = value.replace("`", "")
    value = re.sub(r"\s+", " ", value)
    return value.strip(" -|\n\t")


def _truncate(value: str, limit: int) -> str:
    value = _clean(value)
    return value if len(value) <= limit else value[: limit - 1].rstrip(" ,.;:") + "…"


def _sentence_summary(value: str, limit: int = 210) -> str:
    """Prefer complete sentences over arbitrary character cuts."""
    clean = _clean(value)
    sentences = re.split(r"(?<=[.!?])\s+", clean)
    selected: list[str] = []
    for sentence in sentences:
        candidate = " ".join(selected + [sentence]).strip()
        if len(candidate) > limit:
            break
        selected.append(sentence)
    if selected:
        return " ".join(selected)
    return _truncate(clean, limit)


def _sections(markdown: str) -> dict[str, str]:
    result: dict[str, list[str]] = {}
    current = "overview"
    result[current] = []
    for raw in (markdown or "").splitlines():
        match = re.match(r"^#{1,3}\s+(?:\d+[.)]\s*)?(.+)$", raw.strip())
        if match:
            current = _clean(match.group(1)).lower()
            result.setdefault(current, [])
        else:
            result.setdefault(current, []).append(raw)
    return {key: "\n".join(value).strip() for key, value in result.items()}


def _section(sections: dict[str, str], *needles: str) -> str:
    for key, value in sections.items():
        if any(needle.lower() in key for needle in needles):
            return value
    return ""


def _bullets(text: str, limit: int = 4, max_chars: int = 150) -> list[str]:
    found: list[str] = []
    for raw in (text or "").splitlines():
        match = re.match(r"^\s*(?:[-*]|\d+[.)])\s+(.+)$", raw)
        if not match:
            continue
        value = _truncate(match.group(1), max_chars)
        if value and value not in found:
            found.append(value)
        if len(found) == limit:
            break
    if found:
        return found
    prose = _clean(text)
    sentences = re.split(r"(?<=[.!?])\s+", prose)
    return [_truncate(sentence, max_chars) for sentence in sentences[:limit] if sentence]


def _extract(pattern: str, text: str, default: str = "—") -> str:
    match = re.search(pattern, text or "", re.IGNORECASE | re.MULTILINE | re.DOTALL)
    return _truncate(match.group(1), 80) if match else default


def _label_value(text: str, *labels: str, default: str = "—", limit: int = 70) -> str:
    """Extract the value after a specific Markdown label on the same line."""
    for raw in (text or "").splitlines():
        clean = _clean(raw)
        lowered = clean.lower()
        for label in labels:
            position = lowered.find(label.lower())
            if position < 0:
                continue
            separator = clean.find(":", position + len(label))
            if separator < 0:
                continue
            value = clean[separator + 1 :].strip(" -*")
            if value:
                return _truncate(value, limit)
    return default


def _currency_value(text: str, *labels: str, default: str = "—", pick_last: bool = False) -> str:
    line = _label_value(text, *labels, default="", limit=240)
    values = re.findall(r"\$[\d,.]+(?:\s*(?:billion|million|B|M))?", line, re.IGNORECASE)
    if not values:
        return default
    value = values[-1] if pick_last else values[0]
    value = re.sub(r"\s+billion\b", "B", value, flags=re.IGNORECASE)
    value = re.sub(r"\s+million\b", "M", value, flags=re.IGNORECASE)
    return value


def _set_bg(slide, color=INK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color=WHITE,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return shape


def _rect(slide, x, y, w, h, color=CARD, radius=True, line=None):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line or color
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def _line(slide, x, y, w, color=LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.018))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _chrome(slide, number: int, label: str, title: str, subtitle: str = ""):
    _set_bg(slide)
    _add_text(slide, label.upper(), 0.72, 0.38, 3.8, 0.25, size=9, color=TEAL, bold=True)
    _add_text(slide, title, 0.72, 0.72, 11.8, 0.58, size=28, bold=True)
    if subtitle:
        _add_text(slide, subtitle, 0.74, 1.35, 11.4, 0.42, size=12, color=MUTED)
    _line(slide, 0.72, 6.98, 11.9)
    _add_text(slide, f"{number:02d}", 12.1, 7.08, 0.5, 0.2, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def _card(slide, title: str, body: str, x, y, w, h, accent=TEAL):
    _rect(slide, x, y, w, h)
    _rect(slide, x, y, 0.055, h, accent, radius=False)
    _add_text(slide, title.upper(), x + 0.25, y + 0.23, w - 0.48, 0.25, size=9, color=accent, bold=True)
    _add_text(slide, _truncate(body, 230), x + 0.25, y + 0.63, w - 0.48, h - 0.78, size=14, color=WHITE)


def _bullet_list(slide, items: list[str], x, y, w, h, *, size=15, accent=TEAL):
    if not items:
        items = ["No supporting detail was included in the generated plan."]
    each = h / max(len(items), 1)
    for index, item in enumerate(items):
        cy = y + index * each
        _rect(slide, x, cy + 0.08, 0.09, 0.09, accent, radius=False)
        _add_text(slide, _truncate(item, 170), x + 0.27, cy, w - 0.27, each - 0.03, size=size, color=WHITE)


def _metric(slide, value: str, label: str, x, y, w, accent=TEAL):
    _rect(slide, x, y, w, 1.22)
    value = _truncate(value, 34)
    size = 23 if len(value) <= 12 else 18 if len(value) <= 22 else 14
    _add_text(
        slide,
        value,
        x + 0.15,
        y + 0.14,
        w - 0.3,
        0.62,
        size=size,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    _add_text(slide, label.upper(), x + 0.12, y + 0.79, w - 0.24, 0.22, size=8, color=MUTED, bold=True, align=PP_ALIGN.CENTER)


class PitchDeckGenerator:
    """Create a consistent deck from saved plan data without a second LLM rewrite."""

    def generate(self, plan_data: dict | str) -> bytes:
        if isinstance(plan_data, str):
            plan_data = {"full_plan": plan_data}

        full_plan = plan_data.get("full_plan") or ""
        intelligence = plan_data.get("intelligence") or {}
        identity = intelligence.get("identity") or {}
        metrics_data = intelligence.get("financial_metrics") or {}
        title = _clean(plan_data.get("title") or plan_data.get("idea") or "New Venture")
        score = str(
            (intelligence.get("adjusted_score") or {}).get("adjusted_score")
            or plan_data.get("viability_score")
            or _extract(r"Viability Score:\s*(\d+)", plan_data.get("validation_strategy", ""))
        )
        validation = plan_data.get("validation_strategy") or ""
        financials = plan_data.get("financials") or full_plan
        risks = plan_data.get("risks") or _section(_sections(full_plan), "risk analysis")
        parsed = _sections(full_plan)

        executive = _section(parsed, "executive summary")
        products = _section(parsed, "products & services", "products and services")
        market = _section(parsed, "market analysis")
        model = _section(parsed, "business model")
        marketing = _section(parsed, "marketing & sales", "marketing and sales")
        operations = _section(parsed, "operations plan")
        team = _section(parsed, "management team")
        implementation = _section(parsed, "implementation timeline")
        funding = _section(parsed, "funding requirements")
        revenue_streams = _section(parsed, "revenue streams")
        use_of_funds = _section(parsed, "use of funds")

        verdict = _extract(r"Verdict:\s*([^\n#]+)", validation, "Review")
        market_source = plan_data.get("market_analysis", "") or market
        market_sections = _sections(market_source)
        model_source = plan_data.get("business_model", "") or model
        model_sections = _sections(model_source)
        tam = _currency_value(market_source, "TAM (Total Addressable Market)", default="Not sourced", pick_last=True)
        sam = _currency_value(market_source, "SAM (Serviceable Available Market)", default="Not sourced")
        som = _currency_value(market_source, "SOM (Serviceable Obtainable Market)", default="Not sourced")
        cagr_match = re.search(r"(\d+(?:\.\d+)?%\s*(?:to|–|-)\s*\d+(?:\.\d+)?%)", market_source, re.IGNORECASE)
        cagr = cagr_match.group(1) if cagr_match else _label_value(market_source, "Growth rate", default="Not sourced", limit=24)
        capital_ask = _label_value(
            funding or financials,
            "Recommended Seed Capital",
            "Recommended seed amount",
            default="Not provided",
            limit=28,
        )
        capital_ask = (
            (metrics_data.get("recommended_funding") or {}).get("display")
            or capital_ask
        )
        break_even = _extract(r"Timeline to profitability[^:\n]*:[\s\S]{0,120}?(Month\s+\d+)", financials, "Not proven")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        # 1 — Cover
        slide = prs.slides.add_slide(blank)
        _set_bg(slide)
        _rect(slide, 0, 0, 0.12, 7.5, TEAL, radius=False)
        _add_text(slide, "INVESTOR BRIEF", 0.82, 0.72, 3.5, 0.3, size=10, color=TEAL, bold=True)
        _add_text(slide, title, 0.82, 1.35, 10.9, 1.25, size=42, bold=True)
        one_liner = identity.get("one_liner") or _sentence_summary(executive, 230) or "A structured venture opportunity."
        _add_text(slide, one_liner, 0.86, 2.82, 10.7, 1.05, size=18, color=MUTED)
        _metric(slide, f"{score}/100" if score.isdigit() else score, "viability", 0.86, 5.28, 2.25)
        _metric(slide, verdict.upper(), "decision", 3.32, 5.28, 2.7, CORAL)
        _add_text(slide, "CONFIDENTIAL  •  BUSINESS PLAN", 9.1, 6.95, 3.4, 0.2, size=8, color=MUTED, align=PP_ALIGN.RIGHT)

        # 2 — Decision snapshot
        slide = prs.slides.add_slide(blank)
        _chrome(slide, 2, "Investment snapshot", "The opportunity at a glance")
        _metric(slide, f"{score}/100" if score.isdigit() else score, "viability score", 0.75, 2.0, 2.25)
        _metric(slide, verdict.upper(), "recommendation", 3.17, 2.0, 2.5, CORAL)
        _metric(slide, capital_ask, "capital sought", 5.84, 2.0, 2.65, BLUE)
        _metric(slide, break_even, "claimed break-even", 8.66, 2.0, 2.25)
        timing = _label_value(validation, "Timing", default=f"Market growth is estimated at {cagr}.", limit=145)
        _card(slide, "Why now", _sentence_summary(timing, 145), 0.75, 3.65, 5.72, 2.45, TEAL)
        _card(slide, "Principal concern", (_bullets(risks, 1) or ["Execution risk requires validation."])[0], 6.66, 3.65, 5.72, 2.45, CORAL)

        # 3 — Problem
        slide = prs.slides.add_slide(blank)
        _chrome(slide, 3, "Problem", identity.get("subtitle") or "The customer problem and current alternatives")
        value_prop = _section(model_sections, "value proposition")
        problem_items = _bullets(value_prop, 3) or _bullets(executive, 3)
        _bullet_list(slide, problem_items, 0.82, 2.05, 7.2, 3.9, size=17, accent=CORAL)
        core_tension = problem_items[0] if problem_items else _sentence_summary(plan_data.get("idea", ""), 190)
        _card(slide, "Core tension", core_tension, 8.35, 2.05, 4.0, 3.9, CORAL)

        # 4 — Solution
        slide = prs.slides.add_slide(blank)
        _chrome(slide, 4, "Solution", "A focused product proposition")
        solution_items = _bullets(products, 4) or _bullets(value_prop, 4)
        if not solution_items:
            solution_items = [_sentence_summary(plan_data.get("idea", ""), 170)]
        positions = [(0.75, 2.0), (6.66, 2.0), (0.75, 4.25), (6.66, 4.25)]
        for index, item in enumerate(solution_items[:4]):
            parts = item.split(":", 1)
            card_title = parts[0] if len(parts) == 2 else f"Product 0{index + 1}"
            card_body = parts[1] if len(parts) == 2 else item
            _card(slide, _truncate(card_title, 42), _sentence_summary(card_body, 175), positions[index][0], positions[index][1], 5.72, 1.72, TEAL if index < 2 else BLUE)

        # 5 — Market
        slide = prs.slides.add_slide(blank)
        _chrome(slide, 5, "Market", "Large category, narrow entry wedge")
        for value, label, x in [(tam, "TAM", 0.75), (sam, "SAM", 3.7), (som, "SOM", 6.65), (cagr, "Growth", 9.6)]:
            _metric(slide, value, label, x, 2.0, 2.55)
        target_items = _bullets(_section(market_sections, "target customers"), 2)
        strategy_items = _bullets(marketing, 2) or _bullets(_section(model_sections, "go-to-market"), 2)
        _card(slide, "Initial customer", target_items[0] if target_items else "Target customer evidence is not yet structured.", 0.75, 3.72, 5.72, 2.2, TEAL)
        _card(slide, "Entry strategy", strategy_items[0] if strategy_items else "The initial acquisition strategy requires definition.", 6.66, 3.72, 5.72, 2.2, BLUE)

        # 6 — Business model
        slide = prs.slides.add_slide(blank)
        _chrome(slide, 6, "Business model", "Revenue assumptions and unit economics")
        metric_source = financials or model
        metrics = [
            ((metrics_data.get("arpu") or {}).get("display") or _label_value(metric_source, "ARPU", default="Not provided", limit=20), "ARPU"),
            ((metrics_data.get("gross_margin_rate") or {}).get("display") or _label_value(metric_source, "Gross Margin", default="Not provided", limit=20), "gross margin"),
            ((metrics_data.get("blended_cac") or {}).get("display") or _currency_value(metric_source, "Blended CAC", "Customer Acquisition Cost", default="Not provided"), "blended CAC"),
            ((metrics_data.get("ltv_cac_ratio") or {}).get("display") or _label_value(metric_source, "LTV:CAC ratio", default="Not provided", limit=18), "LTV : CAC"),
        ]
        for index, (value, label) in enumerate(metrics):
            _metric(slide, value, label, 0.75 + index * 2.95, 2.0, 2.55, [TEAL, BLUE, CORAL, TEAL][index])
        _bullet_list(slide, _bullets(revenue_streams or model, 3), 0.82, 3.72, 11.3, 2.15, size=15)

        # 7 — Go-to-market
        slide = prs.slides.add_slide(blank)
        _chrome(slide, 7, "Go-to-market", "Validate demand before scaling inventory")
        gtm = _bullets(marketing, 5)
        _bullet_list(slide, gtm, 0.82, 2.0, 7.0, 4.15, size=16, accent=TEAL)
        _card(slide, "Launch motion", gtm[0] if gtm else "Launch motion is not yet structured.", 8.2, 2.0, 4.15, 2.0, BLUE)
        _card(slide, "Sales expansion", gtm[1] if len(gtm) > 1 else "Expansion channel assumptions require validation.", 8.2, 4.25, 4.15, 1.9, TEAL)

        # 8 — Competition and moat
        slide = prs.slides.add_slide(blank)
        _chrome(slide, 8, "Competition", "Differentiate where incumbents are structurally weak")
        competition_items = _bullets(_section(market_sections, "competitive landscape"), 2)
        moat_items = _bullets(_section(model_sections, "competitive moat"), 1) or _bullets(value_prop, 1)
        _card(slide, "Competitive pressure", competition_items[0] if competition_items else "Competitor evidence is incomplete.", 0.75, 2.0, 3.72, 3.8, CORAL)
        _card(slide, "Market gap", competition_items[1] if len(competition_items) > 1 else "The proposed market gap requires customer validation.", 4.8, 2.0, 3.72, 3.8, BLUE)
        _card(slide, f"Why {title}", moat_items[0] if moat_items else "Differentiation remains a proposed capability.", 8.85, 2.0, 3.72, 3.8, TEAL)

        # 9 — Financial outlook
        slide = prs.slides.add_slide(blank)
        _chrome(slide, 9, "Financial outlook", "The plan assumes rapid scale and early profitability", "Projections are management assumptions and require customer and manufacturing validation.")
        revenue_metrics = [metrics_data.get(f"year_{index}_revenue") or {} for index in range(1, 4)]
        revenues = [
            metric.get("display")
            or _extract(
                rf"Year {index}(?:\s+Projections)?[\s\S]{{0,220}}?"
                rf"Total Revenue:\s*\**\s*\$([0-9,.]+(?:\s*(?:million|M))?)",
                financials,
                "Not provided",
            )
            for index, metric in enumerate(revenue_metrics, 1)
        ]
        revenue_values = [float(metric.get("value") or 0) for metric in revenue_metrics]
        max_revenue = max(revenue_values or [0])
        heights = [0.7 + (value / max_revenue * 2.95 if max_revenue else 0) for value in revenue_values]
        for index, (value, height) in enumerate(zip(revenues, heights)):
            x = 1.25 + index * 3.45
            _rect(slide, x, 6.12 - height, 2.05, height, [BLUE, TEAL, CORAL][index], radius=False)
            _add_text(slide, value, x - 0.1, 5.72 - height, 2.25, 0.35, size=18, bold=True, align=PP_ALIGN.CENTER)
            _add_text(slide, f"YEAR {index + 1}", x, 6.3, 2.05, 0.22, size=9, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        audit_issues = (intelligence.get("consistency_issues") or []) + (intelligence.get("contradiction_issues") or [])
        reality_check = (
            audit_issues[0].get("message")
            or audit_issues[0].get("explanation")
            if audit_issues
            else ((_bullets(risks, 1) or ["Financial projections remain assumptions until validated."])[0])
        )
        _card(slide, "Reality check", reality_check, 10.45, 2.25, 2.15, 3.85, CORAL)

        # 10 — Execution roadmap
        slide = prs.slides.add_slide(blank)
        _chrome(slide, 10, "Roadmap", "De-risk the hardest assumptions in sequence")
        roadmap = _bullets(implementation or operations, 4)
        for index, item in enumerate(roadmap[:4]):
            x = 0.78 + index * 3.0
            _add_text(slide, f"0{index + 1}", x, 2.02, 0.6, 0.4, size=18, color=TEAL, bold=True)
            _line(slide, x, 2.55, 2.55, TEAL)
            _add_text(slide, item, x, 2.85, 2.55, 2.8, size=14)

        # 11 — Risks
        slide = prs.slides.add_slide(blank)
        _chrome(slide, 11, "Risk", "What must be proven before committing full capital")
        risk_items = _bullets(risks, 4)
        for index, item in enumerate(risk_items[:4]):
            y = 1.95 + index * 1.12
            _rect(slide, 0.78, y, 0.45, 0.45, CORAL)
            _add_text(slide, str(index + 1), 0.78, y + 0.08, 0.45, 0.2, size=10, bold=True, align=PP_ALIGN.CENTER)
            _add_text(slide, item, 1.52, y - 0.02, 10.7, 0.75, size=15)

        # 12 — Ask
        slide = prs.slides.add_slide(blank)
        _chrome(slide, 12, "The ask", "Capital tied to measurable de-risking milestones")
        ask = capital_ask
        _add_text(slide, ask, 0.78, 1.9, 5.3, 0.75, size=35, color=TEAL, bold=True)
        _add_text(slide, "RECOMMENDED SEED CAPITAL", 0.82, 2.75, 4.5, 0.25, size=9, color=MUTED, bold=True)
        _bullet_list(slide, _bullets(use_of_funds or funding, 5), 6.15, 1.82, 6.1, 4.55, size=14, accent=BLUE)
        _add_text(slide, "Fund validation—not just launch.", 0.82, 4.15, 4.7, 0.8, size=24, bold=True)
        _add_text(slide, "Release capital in stages as customer demand, delivery feasibility, retention, and acquisition economics are validated.", 0.84, 5.15, 4.8, 1.0, size=14, color=MUTED)

        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()


_pitch_generator = None


def get_pitch_generator() -> PitchDeckGenerator:
    global _pitch_generator
    if _pitch_generator is None:
        _pitch_generator = PitchDeckGenerator()
    return _pitch_generator
